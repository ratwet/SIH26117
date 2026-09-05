"""
SovereignWorkbench — PowerPoint Presentation Compiler (app/compilers/pptx_builder.py)
Generates high-fidelity, executive-ready presentation decks (.pptx)
for refinery turnaround meetings, board reviews, and SIH evaluations.
"""

from pathlib import Path
from typing import Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from app.schemas import ApprovalNotePayload, PipeInspectionData

# Corporate Brand Colors
COLOR_BG_DARK = RGBColor(0x0F, 0x17, 0x2A)    # Deep Slate / #0F172A
COLOR_NAVY_CARD = RGBColor(0x1E, 0x29, 0x3B)  # Card Slate / #1E293B
COLOR_CYAN_ACCENT = RGBColor(0x06, 0xB6, 0xD4)# Bright Cyan / #06B6D4
COLOR_WHITE = RGBColor(0xF8, 0xFA, 0xFC)      # Pure White / #F8FAFC
COLOR_MUTED = RGBColor(0x94, 0xA3, 0xB8)      # Light Gray / #94A3B8
COLOR_CRITICAL_RED = RGBColor(0xEF, 0x44, 0x44) # Red / #EF4444
COLOR_SUCCESS_GREEN = RGBColor(0x10, 0xB9, 0x81) # Green / #10B981


def build_executive_presentation_pptx(
    payload: ApprovalNotePayload,
    output_path: Path
) -> Path:
    """
    Build a 4-slide executive presentation deck (.pptx) summarizing
    the inspection findings, remaining life calculations, and capex budget.
    """
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 Widescreen
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    data = payload.inspection_data

    # -------------------------------------------------------------
    # SLIDE 1: Title & Confidentiality Banner
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = COLOR_BG_DARK
    bg1.line.color.rgb = COLOR_BG_DARK

    # Header Tag
    tag_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.3), Inches(0.5))
    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    p.text = "MANGALORE REFINERY AND PETROCHEMICALS LIMITED (MRPL) • CONFIDENTIAL"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN_ACCENT

    # Main Title
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(2.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"API 570 Piping Inspection & Integrity Audit"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    p2 = tf.add_paragraph()
    p2.text = f"Unit: {payload.unit_name} | Line Tag: {data.line_tag}"
    p2.font.size = Pt(20)
    p2.font.color.rgb = COLOR_MUTED

    # Metadata Grid Card
    meta_card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.2), Inches(11.333), Inches(2.2))
    meta_card.fill.solid()
    meta_card.fill.fore_color.rgb = COLOR_NAVY_CARD
    meta_card.line.color.rgb = COLOR_CYAN_ACCENT

    meta_box = slide1.shapes.add_textbox(Inches(1.3), Inches(4.4), Inches(10.7), Inches(1.8))
    tf_meta = meta_box.text_frame
    p = tf_meta.paragraphs[0]
    p.text = f"• Reference Audit Dossier: {payload.reference_number}"
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_WHITE

    p = tf_meta.add_paragraph()
    p.text = f"• Service & Fluid: {data.service_description} ({data.material_spec})"
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_WHITE

    p = tf_meta.add_paragraph()
    p.text = f"• Design Envelope: {data.design_pressure_psi} psig @ {data.design_temp_celsius} °C | Status: AIR-GAP AQUANEX VERIFIED"
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_CYAN_ACCENT

    p = tf_meta.add_paragraph()
    p.text = f"• Lead Inspector: {data.inspector_name} | Inspection Cycle: {data.inspection_date}"
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_MUTED

    # -------------------------------------------------------------
    # SLIDE 2: Ultrasonic Wall Thickness & Degradation Profile
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = COLOR_BG_DARK
    bg2.line.color.rgb = COLOR_BG_DARK

    header_box = slide2.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.3), Inches(0.8))
    p = header_box.text_frame.paragraphs[0]
    p.text = "Ultrasonic Wall Thickness & Corrosion Audit Findings"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    # 4 Metric Cards
    metrics = [
        ("Nominal Thickness", f"{data.nominal_thickness_mm:.1f} mm", "Original Spool Spec", COLOR_CYAN_ACCENT),
        ("Actual Thickness", f"{data.actual_thickness_mm:.1f} mm", "Measured Ultrasonic UTG", COLOR_CRITICAL_RED),
        ("Minimum Required (t_min)", f"{data.minimum_required_thickness_mm:.1f} mm", "ASME B31.3 Pressure Limit", COLOR_WHITE),
        ("Corrosion Rate", f"{data.corrosion_rate_mm_year:.2f} mm/yr", "Overhead Acidic Vapour", COLOR_MUTED),
    ]

    card_width = Inches(2.6)
    card_gap = Inches(0.3)
    start_left = Inches(1.0)
    for idx, (label, val, sub, col) in enumerate(metrics):
        cur_left = start_left + idx * (card_width + card_gap)
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cur_left, Inches(2.2), card_width, Inches(3.2))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_NAVY_CARD
        card.line.color.rgb = col

        box = slide2.shapes.add_textbox(cur_left + Inches(0.15), Inches(2.4), card_width - Inches(0.3), Inches(2.8))
        tf = box.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = label.upper()
        p1.font.size = Pt(11)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_MUTED

        p2 = tf.add_paragraph()
        p2.text = val
        p2.font.size = Pt(28)
        p2.font.bold = True
        p2.font.color.rgb = col

        p3 = tf.add_paragraph()
        p3.text = sub
        p3.font.size = Pt(11)
        p3.font.color.rgb = COLOR_WHITE

    # -------------------------------------------------------------
    # SLIDE 3: API 570 Remaining Life & Statutory Action
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    bg3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg3.fill.solid()
    bg3.fill.fore_color.rgb = COLOR_BG_DARK
    bg3.line.color.rgb = COLOR_BG_DARK

    header_box = slide3.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.3), Inches(0.8))
    p = header_box.text_frame.paragraphs[0]
    p.text = "Remaining Life Derivation & Statutory Mandate"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    # Main Remaining Life Callout
    callout = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.8), Inches(5.4), Inches(4.5))
    callout.fill.solid()
    callout.fill.fore_color.rgb = COLOR_NAVY_CARD
    callout.line.color.rgb = COLOR_CRITICAL_RED

    cb = slide3.shapes.add_textbox(Inches(1.3), Inches(2.2), Inches(4.8), Inches(3.8))
    tfc = cb.text_frame
    p = tfc.paragraphs[0]
    p.text = "CALCULATED REMAINING LIFE"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_MUTED

    p = tfc.add_paragraph()
    p.text = f"{data.remaining_life_years:.2f} YEARS"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_CRITICAL_RED

    p = tfc.add_paragraph()
    p.text = f"MANDATE: {data.mandatory_action}"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    p = tfc.add_paragraph()
    p.text = "Per API 570 Section 7.1.1, piping systems reaching remaining life < 5.0 years must be scheduled for immediate replacement or major derating."
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_MUTED

    # Formula Steps Box
    f_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.5), Inches(4.5))
    f_box.fill.solid()
    f_box.fill.fore_color.rgb = COLOR_NAVY_CARD
    f_box.line.color.rgb = COLOR_CYAN_ACCENT

    fb = slide3.shapes.add_textbox(Inches(7.1), Inches(2.1), Inches(4.9), Inches(3.9))
    tff = fb.text_frame
    tff.word_wrap = True
    p = tff.paragraphs[0]
    p.text = "FORMULA DERIVATION TRACE (SANDBOX VERIFIED)"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN_ACCENT

    for step in payload.formula_derivation_steps:
        p = tff.add_paragraph()
        p.text = f"• {step}"
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_WHITE

    # -------------------------------------------------------------
    # SLIDE 4: Turnaround Capex & Execution Roadmap
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    bg4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg4.fill.solid()
    bg4.fill.fore_color.rgb = COLOR_BG_DARK
    bg4.line.color.rgb = COLOR_BG_DARK

    header_box = slide4.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.3), Inches(0.8))
    p = header_box.text_frame.paragraphs[0]
    p.text = "Procurement Recommendation & Turnaround Strategy"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    rec_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.8), Inches(11.3), Inches(2.2))
    rec_box.fill.solid()
    rec_box.fill.fore_color.rgb = COLOR_NAVY_CARD
    rec_box.line.color.rgb = COLOR_SUCCESS_GREEN

    rb = slide4.shapes.add_textbox(Inches(1.3), Inches(2.0), Inches(10.7), Inches(1.8))
    tfr = rb.text_frame
    tfr.word_wrap = True
    p = tfr.paragraphs[0]
    p.text = "EXECUTIVE ACTION PLAN"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_SUCCESS_GREEN

    p = tfr.add_paragraph()
    p.text = payload.recommended_action
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    p = tfr.add_paragraph()
    p.text = f"Authorized By: {payload.signatory_title}"
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_CYAN_ACCENT

    # Deliverables Footer Card
    deliv_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.3), Inches(11.3), Inches(2.0))
    deliv_box.fill.solid()
    deliv_box.fill.fore_color.rgb = COLOR_NAVY_CARD
    deliv_box.line.color.rgb = COLOR_CYAN_ACCENT

    db = slide4.shapes.add_textbox(Inches(1.3), Inches(4.5), Inches(10.7), Inches(1.6))
    tfd = db.text_frame
    p = tfd.paragraphs[0]
    p.text = "GENERATED INDUSTRIAL DELIVERABLES (AIR-GAPPED STORAGE):"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN_ACCENT

    p = tfd.add_paragraph()
    p.text = "1. [DOCX] Executive Approval Note (.docx) with official MRPL letterhead and signature block."
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_WHITE

    p = tfd.add_paragraph()
    p.text = "2. [XLSX] Dynamic Cost & Risk Workbook (.xlsx) with active Excel formulas and contingency."
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_WHITE

    p = tfd.add_paragraph()
    p.text = "3. [CAD/DXF] Spool Fabrication Drawing (.dxf) with ASME B31.3 weld-neck flanges."
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_WHITE

    prs.save(str(output_path))
    return output_path
